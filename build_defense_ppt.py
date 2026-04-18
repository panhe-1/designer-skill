from __future__ import annotations

from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


THESIS_PATH = Path(r"C:\Users\林鹿和磐壑\Desktop\论文定稿2.docx")
OUTPUT_PATH = THESIS_PATH.with_name("FD服饰公司短视频营销策略优化研究_终期答辩PPT.pptx")

FONT_CN = "Microsoft YaHei"
FONT_EN = "Aptos"

PALETTE = {
    "bg_light": RGBColor(247, 241, 235),
    "bg_dark": RGBColor(29, 27, 29),
    "text_dark": RGBColor(35, 32, 31),
    "text_mid": RGBColor(90, 82, 76),
    "text_light": RGBColor(244, 239, 233),
    "accent": RGBColor(184, 92, 56),
    "accent_2": RGBColor(121, 135, 107),
    "accent_3": RGBColor(216, 162, 94),
    "line": RGBColor(214, 202, 191),
    "white": RGBColor(255, 255, 255),
    "soft_block": RGBColor(236, 228, 219),
}


def build_presentation() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides = [
        add_cover_slide,
        add_background_slide,
        add_method_slide,
        add_company_slide,
        add_survey_slide,
        add_stp_problem_slide,
        add_4v_problem_slide,
        add_strategy_overview_slide,
        add_stp_strategy_slide,
        add_4v_strategy_slide,
        add_conclusion_slide,
        add_thanks_slide,
    ]

    total = len(slides)
    for index, fn in enumerate(slides, start=1):
        fn(prs, index, total)

    prs.save(str(OUTPUT_PATH))


def hex_line(slide, x: float, y: float, w: float, h: float, fill: RGBColor, transparency: float = 0.0) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.fill.transparency = transparency
    shape.line.fill.background()


def rounded_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: RGBColor,
    line: RGBColor | None = None,
    transparency: float = 0.0,
) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.fill.transparency = transparency
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    if shape.adjustments:
        shape.adjustments[0] = 0.08


def add_textbox(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str = "",
    *,
    font_size: int = 20,
    color: RGBColor | None = None,
    bold: bool = False,
    font_name: str = FONT_CN,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_VERTICAL_ANCHOR = MSO_VERTICAL_ANCHOR.TOP,
    margin: int = 6,
) -> tuple:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Pt(margin)
    tf.margin_right = Pt(margin)
    tf.margin_top = Pt(margin)
    tf.margin_bottom = Pt(margin)
    if text:
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = align
        run = p.runs[0]
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color or PALETTE["text_dark"]
    return box, tf


def add_paragraph(
    tf,
    text: str,
    *,
    font_size: int = 18,
    color: RGBColor | None = None,
    bold: bool = False,
    font_name: str = FONT_CN,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    bullet: bool = False,
    space_after: int = 2,
) -> None:
    p = tf.paragraphs[0] if not tf.text else tf.add_paragraph()
    p.text = text
    p.alignment = align
    p.bullet = bullet
    p.space_after = Pt(space_after)
    if not p.runs:
        p.add_run()
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color or PALETTE["text_dark"]


def set_slide_background(slide, dark: bool = False) -> None:
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = PALETTE["bg_dark" if dark else "bg_light"]


def add_common_chrome(slide, slide_no: int, total: int, title: str, eyebrow: str, dark: bool = False) -> None:
    fg = PALETTE["text_light"] if dark else PALETTE["text_dark"]
    sub = PALETTE["accent_3"] if dark else PALETTE["accent"]
    line = PALETTE["line"] if not dark else PALETTE["accent_2"]

    hex_line(slide, 0.7, 0.48, 0.9, 0.06, sub)
    add_textbox(slide, 0.78, 0.18, 2.8, 0.4, eyebrow, font_size=11, color=sub, bold=True, font_name=FONT_EN)
    add_textbox(slide, 0.72, 0.55, 7.4, 0.6, title, font_size=28, color=fg, bold=True)
    hex_line(slide, 0.72, 1.2, 11.84, 0.02, line)
    add_textbox(
        slide,
        11.7,
        0.2,
        0.9,
        0.35,
        f"{slide_no:02d}/{total:02d}",
        font_size=11,
        color=PALETTE["text_light"] if dark else PALETTE["text_mid"],
        align=PP_ALIGN.RIGHT,
        font_name=FONT_EN,
    )


def add_cover_slide(prs: Presentation, slide_no: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, dark=True)

    hex_line(slide, 0, 0, 13.333, 7.5, PALETTE["bg_dark"])
    hex_line(slide, 9.65, 0, 0.16, 7.5, PALETTE["accent"])
    hex_line(slide, 10.0, 0.65, 2.15, 2.1, PALETTE["accent"], transparency=0.15)
    hex_line(slide, 10.65, 2.15, 1.55, 3.2, PALETTE["accent_2"], transparency=0.12)
    hex_line(slide, 9.9, 5.15, 2.55, 1.2, PALETTE["accent_3"], transparency=0.08)

    rounded_rect(slide, 0.75, 0.75, 2.0, 0.42, PALETTE["accent"], transparency=0.02)
    add_textbox(slide, 0.9, 0.81, 1.7, 0.28, "FINAL DEFENSE", font_size=12, color=PALETTE["text_light"], bold=True, font_name=FONT_EN, align=PP_ALIGN.CENTER, valign=MSO_VERTICAL_ANCHOR.MIDDLE)
    add_textbox(slide, 0.82, 1.55, 8.2, 1.5, "FD服饰公司短视频营销策略优化研究", font_size=30, color=PALETTE["text_light"], bold=True)
    add_textbox(slide, 0.84, 3.02, 6.6, 0.8, "基于 STP + 4V 双理论框架的服饰品牌内容增长研究", font_size=18, color=PALETTE["accent_3"])
    add_textbox(slide, 0.84, 4.05, 7.8, 1.3, "西安欧亚学院  职业教育学院\n市场营销（专升本）", font_size=18, color=PALETTE["text_light"])
    add_textbox(slide, 0.84, 6.0, 3.0, 0.7, "答辩人｜廖子楚", font_size=17, color=PALETTE["text_light"], bold=True)
    add_textbox(slide, 3.2, 6.0, 3.0, 0.7, "指导教师｜巨静文", font_size=17, color=PALETTE["text_light"], bold=True)
    add_textbox(slide, 6.1, 6.0, 2.0, 0.7, f"时间｜{date.today():%Y.%m}", font_size=17, color=PALETTE["text_light"], bold=True, font_name=FONT_EN)
    add_textbox(slide, 10.08, 6.15, 2.0, 0.45, "STP + 4V", font_size=16, color=PALETTE["text_light"], bold=True, font_name=FONT_EN)
    add_textbox(slide, 11.75, 0.2, 0.7, 0.3, f"{slide_no:02d}/{total:02d}", font_size=11, color=PALETTE["accent_3"], align=PP_ALIGN.RIGHT, font_name=FONT_EN)


def add_background_slide(prs: Presentation, slide_no: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_common_chrome(slide, slide_no, total, "研究背景与意义", "RESEARCH BACKGROUND")

    rounded_rect(slide, 0.75, 1.55, 5.1, 4.9, PALETTE["white"], line=PALETTE["line"])
    add_textbox(slide, 1.0, 1.75, 2.0, 0.35, "为什么选这个题目", font_size=13, color=PALETTE["accent"], bold=True)
    _, tf = add_textbox(slide, 1.0, 2.1, 4.45, 3.5)
    for item in [
        "5G 与 AIGC 持续推动内容电商升级，传统流量红利明显收缩。",
        "“短视频 + 直播”已成为服饰品牌的基础配置，但行业信任建设仍明显滞后。",
        "FD 服饰虽已布局抖音、小红书等平台，但尚未形成系统化、可复制的内容打法。",
    ]:
        add_paragraph(tf, item, font_size=18, color=PALETTE["text_mid"], bullet=True, space_after=6)

    rounded_rect(slide, 0.95, 5.35, 4.7, 0.85, PALETTE["soft_block"])
    add_textbox(slide, 1.15, 5.54, 4.3, 0.42, "研究目标：找出 FD 服饰短视频营销的失效根源，并提出可落地的优化路径。", font_size=17, color=PALETTE["text_dark"], bold=True, valign=MSO_VERTICAL_ANCHOR.MIDDLE)

    for (x, y), (value, label) in zip(
        [(6.2, 1.6), (8.45, 1.6), (10.7, 1.6)],
        [("11.21亿", "我国网民规模"), ("95.4%", "短视频使用率"), ("4.8万亿", "2024上半年直播电商交易规模")],
        strict=True,
    ):
        rounded_rect(slide, x, y, 2.0, 1.45, PALETTE["white"], line=PALETTE["line"])
        add_textbox(slide, x + 0.18, y + 0.18, 1.6, 0.45, value, font_size=26, color=PALETTE["accent"], bold=True)
        add_textbox(slide, x + 0.18, y + 0.83, 1.62, 0.4, label, font_size=14, color=PALETTE["text_mid"])

    rounded_rect(slide, 6.2, 3.35, 6.15, 3.1, PALETTE["white"], line=PALETTE["line"])
    add_textbox(slide, 6.45, 3.55, 1.3, 0.34, "研究意义", font_size=13, color=PALETTE["accent"], bold=True)
    add_textbox(slide, 6.45, 3.98, 2.2, 0.35, "理论层面", font_size=20, color=PALETTE["text_dark"], bold=True)
    _, tf_left = add_textbox(slide, 6.45, 4.3, 2.75, 1.6)
    add_paragraph(tf_left, "补充服装品类短视频营销案例研究。", font_size=17, bullet=True, color=PALETTE["text_mid"])
    add_paragraph(tf_left, "验证 STP 与 4V 在细分赛道中的联动应用。", font_size=17, bullet=True, color=PALETTE["text_mid"])
    add_textbox(slide, 9.45, 3.98, 2.2, 0.35, "实践层面", font_size=20, color=PALETTE["text_dark"], bold=True)
    _, tf_right = add_textbox(slide, 9.45, 4.3, 2.55, 1.6)
    add_paragraph(tf_right, "帮助 FD 服饰精准识别内容与转化问题。", font_size=17, bullet=True, color=PALETTE["text_mid"])
    add_paragraph(tf_right, "为同类中小服饰品牌提供可借鉴的优化范式。", font_size=17, bullet=True, color=PALETTE["text_mid"])


def add_method_slide(prs: Presentation, slide_no: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_common_chrome(slide, slide_no, total, "研究思路与方法", "RESEARCH DESIGN")

    rounded_rect(slide, 0.78, 1.58, 7.3, 4.9, PALETTE["white"], line=PALETTE["line"])
    add_textbox(slide, 1.0, 1.8, 1.6, 0.35, "研究逻辑", font_size=13, color=PALETTE["accent"], bold=True)
    stages = [
        ("提出问题", "行业背景 + 企业困境"),
        ("理论构建", "STP + 4V 双理论框架"),
        ("现状诊断", "账号监测 + 问卷数据"),
        ("策略优化", "从方向校准到内容落地"),
    ]
    start_x = 1.02
    for idx, (title, desc) in enumerate(stages):
        x = start_x + idx * 1.72
        rounded_rect(slide, x, 2.4, 1.45, 1.2, PALETTE["soft_block"], line=PALETTE["line"])
        add_textbox(slide, x + 0.12, 2.57, 1.18, 0.32, f"0{idx + 1}", font_size=13, color=PALETTE["accent"], bold=True, font_name=FONT_EN)
        add_textbox(slide, x + 0.12, 2.92, 1.15, 0.45, title, font_size=18, color=PALETTE["text_dark"], bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.08, 3.42, 1.25, 0.52, desc, font_size=13, color=PALETTE["text_mid"], align=PP_ALIGN.CENTER)
        if idx < len(stages) - 1:
            hex_line(slide, x + 1.45, 2.96, 0.28, 0.03, PALETTE["accent_2"])

    rounded_rect(slide, 1.0, 4.55, 6.7, 1.45, PALETTE["soft_block"])
    add_textbox(slide, 1.2, 4.78, 2.1, 0.34, "核心研究问题", font_size=18, color=PALETTE["text_dark"], bold=True)
    _, tf = add_textbox(slide, 1.28, 5.15, 6.15, 0.72)
    add_paragraph(tf, "FD 服饰短视频营销的问题究竟出在战略层还是执行层？", font_size=17, bullet=True, color=PALETTE["text_mid"])
    add_paragraph(tf, "如何通过 STP 校准方向，并用 4V 重构内容与运营动作？", font_size=17, bullet=True, color=PALETTE["text_mid"])

    rounded_rect(slide, 8.4, 1.58, 4.2, 2.0, PALETTE["white"], line=PALETTE["line"])
    add_textbox(slide, 8.65, 1.8, 1.4, 0.35, "研究方法", font_size=13, color=PALETTE["accent"], bold=True)
    add_textbox(slide, 8.68, 2.18, 1.6, 0.4, "文献分析法", font_size=21, color=PALETTE["text_dark"], bold=True)
    add_textbox(slide, 8.68, 2.56, 3.45, 0.55, "梳理短视频营销、服装行业数字化营销与 STP、4V 理论研究。", font_size=16, color=PALETTE["text_mid"])
    add_textbox(slide, 8.68, 3.02, 1.9, 0.4, "问卷调查法", font_size=21, color=PALETTE["text_dark"], bold=True)
    add_textbox(slide, 8.68, 3.4, 3.45, 0.55, "面向时尚穿搭社群、粉丝群与线下门店周边，共回收 412 份有效问卷。", font_size=16, color=PALETTE["text_mid"])

    rounded_rect(slide, 8.4, 3.95, 4.2, 2.53, PALETTE["bg_dark"])
    add_textbox(slide, 8.7, 4.18, 1.8, 0.32, "理论框架", font_size=13, color=PALETTE["accent_3"], bold=True)
    add_textbox(slide, 8.72, 4.6, 1.45, 0.5, "STP", font_size=28, color=PALETTE["text_light"], bold=True, font_name=FONT_EN)
    add_textbox(slide, 10.1, 4.72, 2.1, 0.3, "市场细分 / 目标市场 / 市场定位", font_size=14, color=PALETTE["text_light"])
    add_textbox(slide, 8.72, 5.2, 1.45, 0.5, "4V", font_size=28, color=PALETTE["text_light"], bold=True, font_name=FONT_EN)
    add_textbox(slide, 10.1, 5.26, 2.1, 0.5, "差异化 / 功能化 / 附加价值 / 共鸣", font_size=14, color=PALETTE["text_light"])
    add_textbox(slide, 8.72, 5.95, 3.55, 0.35, "用 STP 找准方向，用 4V 优化内容执行与用户价值传递。", font_size=15, color=PALETTE["accent_3"], bold=True)


def add_company_slide(prs: Presentation, slide_no: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_common_chrome(slide, slide_no, total, "研究对象与营销现状", "CASE CONTEXT")

    rounded_rect(slide, 0.78, 1.58, 4.4, 4.9, PALETTE["white"], line=PALETTE["line"])
    add_textbox(slide, 1.0, 1.82, 1.55, 0.34, "FD 服饰概况", font_size=13, color=PALETTE["accent"], bold=True)
    add_textbox(slide, 1.0, 2.18, 3.5, 0.4, "本土时尚女装品牌", font_size=24, color=PALETTE["text_dark"], bold=True)
    _, tf = add_textbox(slide, 1.0, 2.68, 3.75, 2.2)
    add_paragraph(tf, "2013 年成立，围绕都市年轻女性的通勤、休闲与约会场景布局产品。", font_size=18, bullet=True, color=PALETTE["text_mid"])
    add_paragraph(tf, "价格带集中在 150–500 元，定位“舒适、时尚、高性价比”。", font_size=18, bullet=True, color=PALETTE["text_mid"])
    add_paragraph(tf, "2021 年起布局抖音、小红书等平台，希望借内容营销获取新增量。", font_size=18, bullet=True, color=PALETTE["text_mid"])
    rounded_rect(slide, 1.0, 5.05, 3.95, 0.95, PALETTE["soft_block"])
    add_textbox(slide, 1.18, 5.24, 3.55, 0.46, "现实瓶颈：流量成本抬升，但短视频未形成系统打法，投入产出不成正比。", font_size=16, color=PALETTE["text_dark"], bold=True)

    rounded_rect(slide, 5.45, 1.58, 7.15, 4.9, PALETTE["white"], line=PALETTE["line"])
    add_textbox(slide, 5.7, 1.82, 2.2, 0.34, "现有短视频内容结构", font_size=13, color=PALETTE["accent"], bold=True)
    ratios = [("产品上新展示", 0.60, PALETTE["accent"]), ("日常 vlog / 直播切片", 0.30, PALETTE["accent_2"]), ("穿搭教程", 0.10, PALETTE["accent_3"])]
    bar_x = 5.95
    bar_y = 2.45
    total_w = 5.9
    cursor = bar_x
    for label, share, color in ratios:
        seg_w = total_w * share
        rounded_rect(slide, cursor, bar_y, seg_w, 0.75, color)
        add_textbox(slide, cursor + 0.06, bar_y + 0.15, seg_w - 0.12, 0.46, f"{int(share * 100)}%", font_size=22, color=PALETTE["white"], bold=True, align=PP_ALIGN.CENTER)
        cursor += seg_w
    for idx, (label, share, color) in enumerate(ratios):
        y = 3.45 + idx * 0.65
        hex_line(slide, 5.96, y + 0.11, 0.16, 0.16, color)
        add_textbox(slide, 6.18, y, 2.8, 0.3, label, font_size=17, color=PALETTE["text_dark"], bold=True)
        add_textbox(slide, 10.55, y, 1.2, 0.3, f"{int(share * 100)}%", font_size=17, color=PALETTE["text_mid"], align=PP_ALIGN.RIGHT, font_name=FONT_EN)
    rounded_rect(slide, 5.95, 5.35, 5.95, 0.82, PALETTE["bg_dark"])
    add_textbox(slide, 6.18, 5.53, 5.55, 0.36, "内容形式单一、发布节奏混乱、账号风格摇摆，难以建立稳定认知。", font_size=17, color=PALETTE["text_light"], bold=True)


def add_survey_slide(prs: Presentation, slide_no: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_common_chrome(slide, slide_no, total, "问卷调查与用户画像", "SURVEY INSIGHTS")

    rounded_rect(slide, 0.78, 1.6, 3.2, 4.85, PALETTE["bg_dark"])
    add_textbox(slide, 1.05, 1.88, 1.8, 0.3, "有效样本", font_size=13, color=PALETTE["accent_3"], bold=True)
    add_textbox(slide, 1.02, 2.2, 2.2, 0.9, "412", font_size=46, color=PALETTE["text_light"], bold=True, font_name=FONT_EN)
    add_textbox(slide, 1.05, 3.05, 2.2, 0.3, "份消费者问卷", font_size=18, color=PALETTE["text_light"])
    add_textbox(slide, 1.05, 3.7, 2.2, 0.35, "投放渠道", font_size=18, color=PALETTE["accent_3"], bold=True)
    _, tf = add_textbox(slide, 1.05, 4.08, 2.45, 1.7)
    add_paragraph(tf, "线上：时尚穿搭群、FD 品牌粉丝群", font_size=17, bullet=True, color=PALETTE["text_light"])
    add_paragraph(tf, "线下：FD 门店周边", font_size=17, bullet=True, color=PALETTE["text_light"])
    add_paragraph(tf, "维度：习惯、认知、满意度、购买意愿", font_size=17, bullet=True, color=PALETTE["text_light"])

    rounded_rect(slide, 4.25, 1.6, 8.35, 4.85, PALETTE["white"], line=PALETTE["line"])
    add_textbox(slide, 4.52, 1.86, 1.8, 0.3, "核心画像", font_size=13, color=PALETTE["accent"], bold=True)
    metrics = [("女性占比", 92, "92%"), ("23–30 岁", 65, "65%"), ("一二线城市", 58, "58%"), ("月可支配收入 3000–8000 元", 70, "70%"), ("日均刷短视频 1 小时以上", 82, "82%")]
    for idx, (label, pct, text) in enumerate(metrics):
        y = 2.35 + idx * 0.68
        add_textbox(slide, 4.58, y, 2.7, 0.28, label, font_size=17, color=PALETTE["text_dark"], bold=idx < 2)
        hex_line(slide, 7.15, y + 0.11, 3.2, 0.12, PALETTE["soft_block"])
        hex_line(slide, 7.15, y + 0.11, 3.2 * pct / 100, 0.12, PALETTE["accent" if idx % 2 == 0 else "accent_2"])
        add_textbox(slide, 10.55, y - 0.04, 1.2, 0.3, text, font_size=18, color=PALETTE["text_mid"], align=PP_ALIGN.RIGHT, font_name=FONT_EN)

    rounded_rect(slide, 4.58, 5.15, 7.45, 0.88, PALETTE["soft_block"])
    add_textbox(slide, 4.82, 5.36, 7.0, 0.36, "结论：核心受众是追求高性价比通勤穿搭，同时关注情绪价值与品牌调性的年轻女性。", font_size=17, color=PALETTE["text_dark"], bold=True)


def add_stp_problem_slide(prs: Presentation, slide_no: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_common_chrome(slide, slide_no, total, "问题诊断一：战略层的 STP 失准", "STP DIAGNOSIS")

    cards = [
        ("市场细分过于粗放", "“18–45 岁都市年轻女性”过于宽泛，缺少基于需求与场景的分层。"),
        ("目标市场定位模糊", "核心客群与次级客群不清晰，营销资源分散，内容锚点频繁摇摆。"),
        ("品牌定位落地脱节", "“可持续时尚”未进入内容体系，缺少能持续强化认知的系列化 IP。"),
    ]
    for idx, (title, desc) in enumerate(cards):
        x = 0.86 + idx * 4.15
        rounded_rect(slide, x, 1.85, 3.65, 3.85, PALETTE["white"], line=PALETTE["line"])
        hex_line(slide, x, 1.85, 3.65, 0.14, [PALETTE["accent"], PALETTE["accent_2"], PALETTE["accent_3"]][idx])
        add_textbox(slide, x + 0.2, 2.18, 0.7, 0.3, f"0{idx + 1}", font_size=13, color=PALETTE["text_mid"], bold=True, font_name=FONT_EN)
        add_textbox(slide, x + 0.2, 2.5, 3.05, 0.65, title, font_size=23, color=PALETTE["text_dark"], bold=True)
        add_textbox(slide, x + 0.2, 3.38, 3.0, 1.55, desc, font_size=18, color=PALETTE["text_mid"])

    rounded_rect(slide, 1.0, 6.0, 11.3, 0.62, PALETTE["bg_dark"])
    add_textbox(slide, 1.25, 6.13, 10.8, 0.28, "结果：内容失焦、资源分散、账号价值感不清，用户难以快速建立品牌认知。", font_size=18, color=PALETTE["text_light"], bold=True)


def add_4v_problem_slide(prs: Presentation, slide_no: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_common_chrome(slide, slide_no, total, "问题诊断二：执行层的 4V 断裂", "4V DIAGNOSIS")

    items = [
        ("差异化不足", "内容同质化严重，黄金 3 秒吸引力弱，受众定位也不够精准。"),
        ("功能化不足", "缺少穿搭教程、面料知识与清晰 CTA，内容没能真正解决用户问题。"),
        ("附加价值不足", "品牌故事、环保理念与会员权益传递不足，用户只看到商品而非价值。"),
        ("共鸣不足", "缺少剧情化表达、话题互动与 UGC 机制，难以形成长期情感连接。"),
    ]
    positions = [(0.88, 1.85), (6.72, 1.85), (0.88, 4.05), (6.72, 4.05)]
    colors = [PALETTE["accent"], PALETTE["accent_2"], PALETTE["accent_3"], PALETTE["bg_dark"]]
    for (x, y), (title, desc), color in zip(positions, items, colors, strict=True):
        rounded_rect(slide, x, y, 5.7, 1.75, PALETTE["white"], line=PALETTE["line"])
        hex_line(slide, x, y, 0.18, 1.75, color)
        add_textbox(slide, x + 0.34, y + 0.22, 1.8, 0.35, title, font_size=22, color=PALETTE["text_dark"], bold=True)
        add_textbox(slide, x + 0.34, y + 0.65, 5.0, 0.72, desc, font_size=17, color=PALETTE["text_mid"])

    rounded_rect(slide, 1.0, 6.2, 11.25, 0.42, PALETTE["soft_block"])
    add_textbox(slide, 1.22, 6.24, 10.8, 0.26, "核心矛盾：品牌把短视频做成了“展示渠道”，却没有做成“用户价值载体”。", font_size=17, color=PALETTE["text_dark"], bold=True, align=PP_ALIGN.CENTER)


def add_strategy_overview_slide(prs: Presentation, slide_no: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_common_chrome(slide, slide_no, total, "优化逻辑：先校准方向，再做细内容", "STRATEGY OVERVIEW")

    rounded_rect(slide, 0.95, 2.0, 4.9, 3.4, PALETTE["white"], line=PALETTE["line"])
    rounded_rect(slide, 7.45, 2.0, 4.9, 3.4, PALETTE["white"], line=PALETTE["line"])
    add_textbox(slide, 1.2, 2.3, 1.3, 0.35, "STP", font_size=30, color=PALETTE["accent"], bold=True, font_name=FONT_EN)
    add_textbox(slide, 1.2, 2.72, 2.7, 0.45, "战略校准", font_size=25, color=PALETTE["text_dark"], bold=True)
    _, left_tf = add_textbox(slide, 1.2, 3.28, 4.1, 1.4)
    add_paragraph(left_tf, "精细化市场细分", font_size=19, bullet=True, color=PALETTE["text_mid"])
    add_paragraph(left_tf, "锁定核心与次级目标市场", font_size=19, bullet=True, color=PALETTE["text_mid"])
    add_paragraph(left_tf, "让品牌定位进入内容结构与账号语言", font_size=19, bullet=True, color=PALETTE["text_mid"])

    add_textbox(slide, 7.72, 2.3, 1.3, 0.35, "4V", font_size=30, color=PALETTE["accent_2"], bold=True, font_name=FONT_EN)
    add_textbox(slide, 7.72, 2.72, 2.7, 0.45, "执行优化", font_size=25, color=PALETTE["text_dark"], bold=True)
    _, right_tf = add_textbox(slide, 7.72, 3.28, 4.2, 1.5)
    add_paragraph(right_tf, "做出差异化内容表达", font_size=19, bullet=True, color=PALETTE["text_mid"])
    add_paragraph(right_tf, "增强内容实用性与转化效率", font_size=19, bullet=True, color=PALETTE["text_mid"])
    add_paragraph(right_tf, "补齐附加价值与情感共鸣", font_size=19, bullet=True, color=PALETTE["text_mid"])

    hex_line(slide, 5.95, 3.55, 1.0, 0.06, PALETTE["accent"])
    hex_line(slide, 6.8, 3.42, 0.16, 0.32, PALETTE["accent"])
    rounded_rect(slide, 4.9, 5.7, 3.55, 0.62, PALETTE["bg_dark"])
    add_textbox(slide, 5.15, 5.83, 3.05, 0.26, "预期结果：提升识别度、转化率、复购意愿与品牌资产沉淀", font_size=16, color=PALETTE["text_light"], bold=True, align=PP_ALIGN.CENTER)


def add_stp_strategy_slide(prs: Presentation, slide_no: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_common_chrome(slide, slide_no, total, "基于 STP 的优化策略", "STP SOLUTIONS")

    rounded_rect(slide, 0.82, 1.62, 5.55, 4.92, PALETTE["white"], line=PALETTE["line"])
    add_textbox(slide, 1.05, 1.87, 1.8, 0.32, "目标客群分层", font_size=13, color=PALETTE["accent"], bold=True)
    layers = [
        ("核心目标", "精致通勤族", "22–28 岁，关注一衣多穿与职场性价比。"),
        ("次级目标", "悦己体验党", "29–45 岁，更看重品质、设计与情绪价值。"),
        ("机会圈层", "可持续生活家", "关注环保面料、安全性与品牌理念一致性。"),
    ]
    widths = [4.2, 3.4, 2.6]
    colors = [PALETTE["accent"], PALETTE["accent_2"], PALETTE["accent_3"]]
    y = 2.45
    for idx, (tag, title, desc) in enumerate(layers):
        w = widths[idx]
        x = 1.15 + (4.2 - w) / 2
        rounded_rect(slide, x, y + idx * 1.08, w, 0.92, colors[idx])
        add_textbox(slide, x + 0.15, y + 0.1 + idx * 1.08, 0.9, 0.24, tag, font_size=12, color=PALETTE["white"], bold=True)
        add_textbox(slide, x + 0.15, y + 0.35 + idx * 1.08, w - 0.3, 0.22, title, font_size=20, color=PALETTE["white"], bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, 1.05, y + 0.84 + idx * 1.08, 4.95, 0.2, desc, font_size=14, color=PALETTE["text_mid"], align=PP_ALIGN.CENTER)

    rounded_rect(slide, 6.62, 1.62, 5.85, 4.92, PALETTE["white"], line=PALETTE["line"])
    add_textbox(slide, 6.88, 1.87, 1.9, 0.32, "内容结构重组", font_size=13, color=PALETTE["accent"], bold=True)
    ratios = [
        ("场景化通勤教程", 0.40, PALETTE["accent"]),
        ("职场穿搭解决方案", 0.30, PALETTE["accent_2"]),
        ("产品展示", 0.20, PALETTE["accent_3"]),
        ("品牌故事 / 可持续理念", 0.10, PALETTE["bg_dark"]),
    ]
    base_x = 7.0
    base_y = 2.55
    total_h = 3.0
    cursor = base_y
    for label, share, color in ratios:
        h = total_h * share
        rounded_rect(slide, base_x, cursor, 1.2, h, color)
        add_textbox(slide, base_x + 0.08, cursor + h / 2 - 0.14, 1.02, 0.28, f"{int(share * 100)}%", font_size=16, color=PALETTE["white"], bold=True, align=PP_ALIGN.CENTER, valign=MSO_VERTICAL_ANCHOR.MIDDLE)
        cursor += h
    for idx, (label, share, color) in enumerate(ratios):
        y = 2.55 + idx * 0.74
        hex_line(slide, 8.45, y + 0.11, 0.16, 0.16, color)
        add_textbox(slide, 8.7, y, 3.05, 0.28, label, font_size=17, color=PALETTE["text_dark"], bold=True)

    rounded_rect(slide, 7.0, 5.35, 5.1, 0.78, PALETTE["soft_block"])
    add_textbox(slide, 7.22, 5.55, 4.7, 0.3, "账号锚点：围绕“通勤穿搭指南”等系列 IP，强化可持续时尚解决方案。", font_size=16, color=PALETTE["text_dark"], bold=True)


def add_4v_strategy_slide(prs: Presentation, slide_no: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_common_chrome(slide, slide_no, total, "基于 4V 的内容与运营策略", "4V SOLUTIONS")

    items = [
        ("差异化", "真人出镜 + 痛点前置", "围绕通勤、约会、休闲场景设计系列栏目，抢占用户“黄金 3 秒”。", PALETTE["accent"]),
        ("功能化", "做教程而非只卖货", "输出一衣多穿、版型选择、面料打理等实用内容，并补齐 CTA 与直播转化钩子。", PALETTE["accent_2"]),
        ("附加价值", "把品牌理念说清楚", "用品牌故事、设计师幕后、环保实践与会员权益建立长期竞争壁垒。", PALETTE["accent_3"]),
        ("共鸣", "从单向输出走向共创", "通过微剧情、话题挑战、UGC 征集与评论互动建立情感与社群连接。", PALETTE["bg_dark"]),
    ]
    positions = [(0.9, 1.82), (6.68, 1.82), (0.9, 4.08), (6.68, 4.08)]
    for (x, y), (tag, title, desc, color) in zip(positions, items, strict=True):
        rounded_rect(slide, x, y, 5.45, 1.85, PALETTE["white"], line=PALETTE["line"])
        hex_line(slide, x, y, 5.45, 0.12, color)
        add_textbox(slide, x + 0.22, y + 0.22, 0.95, 0.24, tag, font_size=13, color=color, bold=True)
        add_textbox(slide, x + 0.22, y + 0.5, 1.95, 0.32, title, font_size=22, color=PALETTE["text_dark"], bold=True)
        add_textbox(slide, x + 0.22, y + 0.92, 4.92, 0.58, desc, font_size=17, color=PALETTE["text_mid"])

    rounded_rect(slide, 1.0, 6.16, 11.2, 0.42, PALETTE["bg_dark"])
    add_textbox(slide, 1.3, 6.21, 10.6, 0.24, "最终目标：让 FD 服饰从“卖衣服”转向“用内容驱动消费、用品牌沉淀复购”。", font_size=17, color=PALETTE["text_light"], bold=True, align=PP_ALIGN.CENTER)


def add_conclusion_slide(prs: Presentation, slide_no: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_common_chrome(slide, slide_no, total, "研究结论、局限与展望", "CONCLUSION")

    rounded_rect(slide, 0.8, 1.6, 6.2, 4.95, PALETTE["white"], line=PALETTE["line"])
    add_textbox(slide, 1.05, 1.86, 1.8, 0.3, "研究结论", font_size=13, color=PALETTE["accent"], bold=True)
    conclusion_cards = [
        ("根源判断", "问题不在“发得不够多”，而在战略失焦与执行偏差并存。"),
        ("解决逻辑", "先用 STP 重构方向，再用 4V 提升内容质量、价值感与共鸣度。"),
        ("实践价值", "为中小服饰品牌提供了一套从用户价值出发的短视频优化路径。"),
    ]
    for idx, (title, desc) in enumerate(conclusion_cards):
        y = 2.25 + idx * 1.08
        rounded_rect(slide, 1.05, y, 5.65, 0.82, PALETTE["soft_block"])
        add_textbox(slide, 1.28, y + 0.16, 1.2, 0.24, title, font_size=17, color=PALETTE["text_dark"], bold=True)
        add_textbox(slide, 2.18, y + 0.16, 4.2, 0.4, desc, font_size=16, color=PALETTE["text_mid"])

    rounded_rect(slide, 7.3, 1.6, 5.18, 2.15, PALETTE["white"], line=PALETTE["line"])
    add_textbox(slide, 7.56, 1.86, 1.4, 0.3, "研究局限", font_size=13, color=PALETTE["accent"], bold=True)
    _, tf1 = add_textbox(slide, 7.56, 2.22, 4.5, 1.2)
    add_paragraph(tf1, "单案例研究，普适性有限。", font_size=16, bullet=True, color=PALETTE["text_mid"])
    add_paragraph(tf1, "样本以一二线年轻女性为主。", font_size=16, bullet=True, color=PALETTE["text_mid"])
    add_paragraph(tf1, "优化策略尚未经过长期市场实证检验。", font_size=16, bullet=True, color=PALETTE["text_mid"])

    rounded_rect(slide, 7.3, 4.05, 5.18, 2.5, PALETTE["bg_dark"])
    add_textbox(slide, 7.56, 4.3, 1.7, 0.3, "未来展望", font_size=13, color=PALETTE["accent_3"], bold=True)
    _, tf2 = add_textbox(slide, 7.56, 4.66, 4.42, 1.55)
    add_paragraph(tf2, "借助 AI 提升内容洞察、文案生成与虚拟试穿效率。", font_size=16, bullet=True, color=PALETTE["text_light"])
    add_paragraph(tf2, "建立“头部 + 腰部 + 素人”金字塔达人合作矩阵。", font_size=16, bullet=True, color=PALETTE["text_light"])
    add_paragraph(tf2, "优化组织架构与考核体系，形成内容驱动的长期增长闭环。", font_size=16, bullet=True, color=PALETTE["text_light"])


def add_thanks_slide(prs: Presentation, slide_no: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, dark=True)
    hex_line(slide, 0, 0, 13.333, 7.5, PALETTE["bg_dark"])
    hex_line(slide, 0.9, 1.1, 2.6, 0.08, PALETTE["accent"])
    hex_line(slide, 9.85, 0.9, 2.0, 0.12, PALETTE["accent_3"], transparency=0.1)
    hex_line(slide, 10.6, 1.6, 1.1, 3.8, PALETTE["accent"], transparency=0.12)
    hex_line(slide, 9.2, 4.95, 2.7, 1.05, PALETTE["accent_2"], transparency=0.1)
    add_textbox(slide, 0.95, 1.45, 5.2, 0.5, "THANK YOU", font_size=13, color=PALETTE["accent_3"], bold=True, font_name=FONT_EN)
    add_textbox(slide, 0.95, 2.0, 4.8, 1.0, "谢谢聆听", font_size=34, color=PALETTE["text_light"], bold=True)
    add_textbox(slide, 0.95, 3.05, 5.6, 0.55, "恳请各位老师批评指正", font_size=20, color=PALETTE["text_light"])
    add_textbox(slide, 0.95, 4.55, 5.6, 0.7, "FD服饰公司短视频营销策略优化研究", font_size=19, color=PALETTE["accent_3"], bold=True)
    add_textbox(slide, 0.95, 5.5, 4.6, 0.55, "答辩人｜廖子楚    指导教师｜巨静文", font_size=17, color=PALETTE["text_light"], bold=True)
    add_textbox(slide, 11.75, 0.2, 0.7, 0.3, f"{slide_no:02d}/{total:02d}", font_size=11, color=PALETTE["accent_3"], align=PP_ALIGN.RIGHT, font_name=FONT_EN)


if __name__ == "__main__":
    build_presentation()
