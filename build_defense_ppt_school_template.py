from __future__ import annotations

from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
MEDIA_DIR = BASE_DIR / "template_media"
OUTPUT_PATH = Path.home() / "Desktop" / "FD服饰公司短视频营销策略优化研究_终期答辩PPT_学校模板增强版.pptx"

FONT_CN = "Microsoft YaHei"
FONT_EN = "Aptos"

THESIS_TITLE = "FD服饰公司短视频营销策略优化研究"
SUBTITLE = "基于 STP 战略与 4V 营销理论的诊断与优化"
AUTHOR = "廖子楚"
ADVISOR = "巨静文"
COLLEGE = "西安欧亚学院 职业教育学院"
MAJOR = "市场营销专业"

PALETTE = {
    "white": RGBColor(255, 255, 255),
    "bg": RGBColor(249, 251, 252),
    "ink": RGBColor(33, 44, 54),
    "muted": RGBColor(96, 111, 123),
    "teal": RGBColor(49, 191, 181),
    "teal_dark": RGBColor(28, 134, 133),
    "teal_deep": RGBColor(17, 91, 102),
    "panel": RGBColor(214, 229, 245),
    "panel_soft": RGBColor(236, 245, 251),
    "line": RGBColor(191, 219, 226),
    "mint": RGBColor(226, 244, 240),
    "gold": RGBColor(216, 168, 88),
    "orange": RGBColor(215, 142, 86),
    "slate": RGBColor(128, 141, 150),
    "shadow": RGBColor(227, 235, 240),
    "navy": RGBColor(54, 73, 93),
}

IMAGES = {
    "cover": MEDIA_DIR / "image1.jpeg",
    "contents": MEDIA_DIR / "image4.jpeg",
    "section1": MEDIA_DIR / "image4.jpeg",
    "section2": MEDIA_DIR / "image6.jpeg",
    "section3": MEDIA_DIR / "image7.jpeg",
    "section4": MEDIA_DIR / "image8.jpeg",
    "section5": MEDIA_DIR / "image5.jpeg",
    "thanks": MEDIA_DIR / "image1.jpeg",
    "logo": MEDIA_DIR / "image10.png",
}

DOMESTIC_REFS_1 = [
    "[1] 侯雪青. 短视频的营销模式及策略研究[J]. 电子商务评论, 2024, 13(2):2298-2303.",
    "[2] 赵婷. 服装类短视频内容营销影响购买意愿的作用机理研究[D]. 东华大学, 2023.",
    "[3] 舒方瑞. 服装市场运营中短视频应用的优势和策略[J]. 上海服饰, 2024(04):8-10.",
    "[4] 汪彦丹. 中国服饰品牌短视频传播策略研究分析: 以太平鸟女装为例[J]. 化纤与纺织技术, 2022, 51(08):109-111.",
    "[5] 姜雯静, 孙虹. 基于 SIPS 模型的服装品牌短视频营销体系构建[J]. 服装设计师, 2024(10):122-128.",
    "[6] 吴金明. 新经济时代的“4V”营销组合[J]. 中国工业经济, 2001(06):70-75.",
    "[7] 严晓冬. 新媒体短视频的特征及其内容创作策略[J]. 西部广播电视, 2023, 44(07):57-59.",
    "[8] 李庆, 秦娟. 数字化营销时代短视频的应用研究[J]. 辽宁省交通高等专科学校学报, 2025, 27(04):28-31.",
    "[9] 于潇, 王鼎立, 白丽, 等. 内容策略视域下的企业号短视频用户参与行为研究[J]. 情报科学, 2023, 41(11):85-93,150.",
    "[10] 程玉田. 4I 营销理论视角下虚拟数字人营销实践的现状、问题及优化对策研究[D]. 上海师范大学, 2023.",
]

DOMESTIC_REFS_2 = [
    "[11] 唐丽琼. JTY 公司短视频营销策略研究[D]. 桂林电子科技大学, 2023.",
    "[12] 程玲玲. 武汉 A 传媒公司短视频业务营销策略优化研究[D]. 华中农业大学, 2023.",
    "[13] 严小林. 短视频平台中商品展示的视觉修辞与消费行为引导机制探究[J]. 人像摄影, 2025(09):262-263.",
    "[14] 珍珍. 说服理论视域下商品导购类短视频对内蒙古地区大学生消费行为影响研究[D]. 内蒙古师范大学, 2024.",
    "[15] 刘雅婷, 陈文晖. “十五五”时期我国纺织服装产业数字化升级路径研究[J]. 价格理论与实践, 2025(12):90-96+290.",
]

FOREIGN_REFS = [
    "[16] Tam Y F, Lung J. Digital marketing strategies for luxury fashion brands: A systematic literature review[J]. International Journal of Information Management Data Insights, 2025, 5(1):100309.",
    "[17] Kanellos N, Karountzos P, Giannakopoulos T N, et al. Digital Marketing Strategies and Profitability in the Agri-Food Industry[J]. Sustainability, 2024, 16(14):5889.",
    "[18] Luo C, Hasan M A N, Ahmad B Z M A, et al. Influence of short video content on consumers' purchase intentions on social media platforms[J]. Scientific Reports, 2025, 15(1):16605.",
    "[19] Zafar A U, Shen J, Shahzad M, Islam T. Social media and sustainable purchasing attitude[J]. Journal of Retailing and Consumer Services, 2021, 63:102751.",
]


def set_background(slide, color: RGBColor = PALETTE["bg"]) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def style_run(run, *, font_size: int, color: RGBColor, bold: bool = False, font_name: str = FONT_CN) -> None:
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    rpr = run._element.get_or_add_rPr()
    rpr.set(qn("a:latin"), font_name)
    rpr.set(qn("a:ea"), font_name)
    rpr.set(qn("a:cs"), font_name)


def add_shape(slide, shape_type, x, y, w, h, fill, *, line=None, transparency: float = 0.0):
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.fill.transparency = transparency
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def add_rect(slide, x, y, w, h, fill, *, line=None, transparency: float = 0.0):
    return add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h, fill, line=line, transparency=transparency)


def add_round_rect(slide, x, y, w, h, fill, *, line=None, transparency: float = 0.0):
    shape = add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        x,
        y,
        w,
        h,
        fill,
        line=line,
        transparency=transparency,
    )
    if shape.adjustments:
        shape.adjustments[0] = 0.08
    return shape


def add_circle(slide, x, y, d, fill, *, line=None, transparency: float = 0.0):
    return add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, x, y, d, d, fill, line=line, transparency=transparency)


def add_textbox(
    slide,
    x,
    y,
    w,
    h,
    text: str = "",
    *,
    font_size: int = 18,
    color: RGBColor = PALETTE["ink"],
    bold: bool = False,
    font_name: str = FONT_CN,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_VERTICAL_ANCHOR = MSO_VERTICAL_ANCHOR.TOP,
    margin: int = 4,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Pt(margin)
    tf.margin_right = Pt(margin)
    tf.margin_top = Pt(margin)
    tf.margin_bottom = Pt(margin)
    if text:
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = align
        if not p.runs:
            p.add_run()
        style_run(p.runs[0], font_size=font_size, color=color, bold=bold, font_name=font_name)
    return box, tf


def add_paragraph(
    tf,
    text: str,
    *,
    font_size: int = 16,
    color: RGBColor = PALETTE["ink"],
    bold: bool = False,
    font_name: str = FONT_CN,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    bullet: bool = False,
    level: int = 0,
    space_after: int = 2,
    line_spacing: float = 1.1,
):
    p = tf.paragraphs[0] if not tf.text else tf.add_paragraph()
    p.text = text
    p.alignment = align
    p.level = level
    p.bullet = bullet
    p.space_after = Pt(space_after)
    p.line_spacing = line_spacing
    if not p.runs:
        p.add_run()
    style_run(p.runs[0], font_size=font_size, color=color, bold=bold, font_name=font_name)
    return p


def add_chip(slide, x, y, w, h, text: str, *, fill, color, bold: bool = True, font_size: int = 12):
    add_round_rect(slide, x, y, w, h, fill)
    add_textbox(
        slide,
        x + 0.04,
        y + 0.02,
        w - 0.08,
        h - 0.04,
        text,
        font_size=font_size,
        color=color,
        bold=bold,
        align=PP_ALIGN.CENTER,
        valign=MSO_VERTICAL_ANCHOR.MIDDLE,
        margin=0,
    )


def add_panel_title(slide, x, y, title: str, *, size: int = 18, accent: bool = True):
    if accent:
        add_rect(slide, x, y + 0.06, 0.08, 0.24, PALETTE["teal"])
        add_textbox(slide, x + 0.14, y, 3.5, 0.34, title, font_size=size, color=PALETTE["ink"], bold=True)
    else:
        add_textbox(slide, x, y, 3.5, 0.34, title, font_size=size, color=PALETTE["ink"], bold=True)


def add_metric_card(slide, x, y, w, h, value: str, label: str, *, value_size: int = 24):
    add_round_rect(slide, x, y, w, h, PALETTE["white"], line=PALETTE["line"])
    add_textbox(slide, x + 0.12, y + 0.14, w - 0.24, 0.34, value, font_size=value_size, color=PALETTE["teal_dark"], bold=True, font_name=FONT_EN)
    add_textbox(slide, x + 0.12, y + 0.54, w - 0.24, 0.38, label, font_size=13, color=PALETTE["muted"])


def add_top_chrome(slide, title: str, slide_no: int, total: int, *, subtitle: str | None = None) -> None:
    set_background(slide)
    add_textbox(slide, 0.26, 0.16, 0.3, 0.34, "•", font_size=26, color=PALETTE["teal_dark"], bold=True, font_name=FONT_EN)
    add_textbox(slide, 0.58, 0.14, 5.7, 0.42, title, font_size=28, color=PALETTE["ink"], bold=True)
    if subtitle:
        add_textbox(slide, 0.62, 0.54, 4.4, 0.24, subtitle, font_size=11, color=PALETTE["muted"], font_name=FONT_EN)
    add_round_rect(slide, 12.42, 0.18, 0.22, 0.12, PALETTE["teal"])
    add_textbox(
        slide,
        11.98,
        0.14,
        0.38,
        0.22,
        f"{slide_no:02d}",
        font_size=11,
        color=PALETTE["muted"],
        font_name=FONT_EN,
        align=PP_ALIGN.RIGHT,
    )
    add_textbox(slide, 12.62, 0.14, 0.34, 0.22, f"/ {total:02d}", font_size=9, color=PALETTE["muted"], font_name=FONT_EN)


def add_footer_bar(slide) -> None:
    add_rect(slide, 0.24, 7.03, 12.45, 0.07, PALETTE["teal"], transparency=0.1)


def add_picture_fill(slide, image: Path, x, y, w, h):
    if image.exists():
        return slide.shapes.add_picture(str(image), Inches(x), Inches(y), Inches(w), Inches(h))
    return None


def add_reference_column(slide, x, y, w, h, title: str, refs: list[str]) -> None:
    add_round_rect(slide, x, y, w, h, PALETTE["white"], line=PALETTE["line"])
    add_panel_title(slide, x + 0.2, y + 0.16, title, size=16)
    _, tf = add_textbox(slide, x + 0.18, y + 0.54, w - 0.36, h - 0.72)
    for ref in refs:
        add_paragraph(tf, ref, font_size=11, color=PALETTE["muted"], space_after=6, line_spacing=1.05)


def build_cover_slide(prs: Presentation, slide_no: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, PALETTE["white"])

    add_rect(slide, 0, 0, 13.333, 7.5, PALETTE["white"])
    add_rect(slide, 0, 0, 7.9, 7.5, PALETTE["teal"], transparency=0.84)
    add_rect(slide, 0, 0, 13.333, 1.0, PALETTE["teal"], transparency=0.94)
    add_rect(slide, 9.95, 1.0, 2.55, 5.55, PALETTE["shadow"])
    add_picture_fill(slide, IMAGES["cover"], 7.72, 1.02, 4.62, 5.55)
    if IMAGES["logo"].exists():
        add_picture_fill(slide, IMAGES["logo"], 0.42, 0.22, 1.38, 0.34)

    add_textbox(slide, 0.58, 1.42, 5.8, 0.55, "本科毕业论文终期答辩", font_size=26, color=PALETTE["teal"], bold=True)
    add_textbox(slide, 0.58, 2.18, 6.05, 1.5, THESIS_TITLE, font_size=33, color=PALETTE["white"], bold=True)
    add_textbox(slide, 0.62, 3.88, 5.8, 0.72, SUBTITLE, font_size=19, color=PALETTE["white"], bold=True)

    add_round_rect(slide, 0.6, 4.9, 5.95, 1.18, PALETTE["white"], transparency=0.08)
    add_textbox(slide, 0.86, 5.12, 4.8, 0.32, f"{COLLEGE}  |  {MAJOR}", font_size=16, color=PALETTE["ink"], bold=True)
    add_textbox(slide, 0.86, 5.54, 1.75, 0.28, f"答辩人：{AUTHOR}", font_size=16, color=PALETTE["ink"], bold=True)
    add_textbox(slide, 2.74, 5.54, 2.15, 0.28, f"指导教师：{ADVISOR}", font_size=16, color=PALETTE["ink"], bold=True)
    add_textbox(slide, 5.04, 5.54, 1.25, 0.28, f"{date.today():%Y年%m月}", font_size=16, color=PALETTE["ink"], bold=True, font_name=FONT_EN)
    add_textbox(slide, 12.1, 0.18, 0.35, 0.22, f"{slide_no:02d}", font_size=10, color=PALETTE["muted"], font_name=FONT_EN)


def build_contents_slide(prs: Presentation, slide_no: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, PALETTE["white"])

    add_picture_fill(slide, IMAGES["contents"], 0.42, 0.42, 3.62, 6.4)
    add_rect(slide, 0.42, 0.42, 3.62, 6.4, PALETTE["teal"], transparency=0.76)
    add_rect(slide, 0.42, 0.42, 1.0, 6.4, PALETTE["white"], transparency=0.26)
    add_textbox(slide, 1.5, 0.96, 2.0, 0.34, "CONTENTS", font_size=18, color=PALETTE["white"], font_name=FONT_EN, bold=True)
    add_textbox(slide, 1.26, 5.08, 2.1, 0.72, "目录", font_size=38, color=PALETTE["white"], bold=True)
    add_textbox(slide, 1.3, 5.9, 2.15, 0.52, "终期答辩结构与汇报重点", font_size=17, color=PALETTE["white"], bold=True)

    add_round_rect(slide, 4.38, 0.62, 8.5, 6.05, PALETTE["panel_soft"], line=PALETTE["line"])
    add_textbox(slide, 12.16, 0.18, 0.35, 0.22, f"{slide_no:02d}", font_size=10, color=PALETTE["muted"], font_name=FONT_EN)

    items = [
        ("01", "研究背景、文献与理论", "研究缘起、国内外研究现状、理论依据"),
        ("02", "研究设计", "研究方法、样本来源、技术路线"),
        ("03", "现状分析与问题诊断", "企业现状、问卷发现、STP 与 4V 问题"),
        ("04", "优化策略", "STP 重构、4V 优化、落地重点"),
        ("05", "研究结论与展望", "主要结论、研究局限、未来方向"),
        ("06", "主要参考文献", "按论文参考文献整理展示"),
    ]
    y = 0.98
    for num, title, desc in items:
        add_round_rect(slide, 4.74, y, 7.75, 0.84, PALETTE["white"], line=PALETTE["line"])
        add_circle(slide, 5.0, y + 0.15, 0.46, PALETTE["teal"], transparency=0.04)
        add_textbox(slide, 5.06, y + 0.22, 0.34, 0.15, num, font_size=13, color=PALETTE["teal_dark"], bold=True, font_name=FONT_EN, align=PP_ALIGN.CENTER)
        add_textbox(slide, 5.68, y + 0.1, 3.7, 0.28, title, font_size=21, color=PALETTE["ink"], bold=True)
        add_textbox(slide, 9.0, y + 0.13, 3.18, 0.28, desc, font_size=15, color=PALETTE["muted"])
        y += 0.9

def build_section_slide(prs: Presentation, slide_no: int, total: int, num: str, title: str, image: Path, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, PALETTE["white"])
    add_rect(slide, 0, 0, 13.333, 7.5, PALETTE["white"])
    add_rect(slide, 0, 0, 8.25, 7.5, PALETTE["panel_soft"])
    add_rect(slide, 0, 5.82, 13.333, 1.68, PALETTE["teal"], transparency=0.93)
    add_picture_fill(slide, image, 8.48, 0.82, 3.72, 5.06)
    add_rect(slide, 8.48, 0.82, 3.72, 5.06, PALETTE["teal"], transparency=0.82)
    add_textbox(slide, 1.18, 1.72, 1.8, 0.78, num, font_size=40, color=PALETTE["teal"], bold=True, font_name=FONT_EN)
    add_round_rect(slide, 1.16, 2.95, 5.25, 0.78, PALETTE["teal"])
    add_textbox(slide, 1.48, 3.1, 4.72, 0.4, title, font_size=28, color=PALETTE["white"], bold=True, valign=MSO_VERTICAL_ANCHOR.MIDDLE)
    add_textbox(slide, 1.22, 4.02, 3.85, 0.3, subtitle, font_size=13, color=PALETTE["muted"], font_name=FONT_EN, bold=True)
    add_textbox(slide, 1.22, 4.5, 4.85, 0.56, "围绕论文核心逻辑展开：从背景与理论出发，进入诊断，再落到优化路径与结论。", font_size=16, color=PALETTE["ink"])
    add_textbox(slide, 12.15, 0.18, 0.35, 0.22, f"{slide_no:02d}", font_size=10, color=PALETTE["muted"], font_name=FONT_EN)


def build_background_slide(prs: Presentation, slide_no: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_chrome(slide, "研究背景与意义", slide_no, total, subtitle="RESEARCH BACKGROUND")
    add_round_rect(slide, 0.6, 1.0, 12.1, 5.88, PALETTE["panel_soft"], line=PALETTE["line"])

    for x, value, label in [
        (0.92, "11.24亿", "截至 2024 年 6 月我国网民规模"),
        (2.98, "97.6%", "短视频用户使用率"),
        (5.04, "4.8万亿", "2024 年上半年直播电商交易规模"),
    ]:
        add_metric_card(slide, x, 1.28, 1.8, 1.18, value, label, value_size=25)

    add_round_rect(slide, 0.86, 2.74, 5.25, 3.55, PALETTE["white"], line=PALETTE["line"])
    add_panel_title(slide, 1.04, 2.94, "研究背景", size=18)
    _, tf_left = add_textbox(slide, 1.0, 3.38, 4.82, 2.34)
    for item in [
        "5G、AIGC 与内容电商深度融合，服饰品牌竞争从“铺量”转向“内容价值”。",
        "短视频与直播已成为服饰品牌获取曝光、互动和转化的重要触点。",
        "FD 服饰已入驻抖音、小红书等平台，但内容仍以产品展示为主，缺少场景化表达。",
        "因此，有必要从理论与实践两端诊断其短视频营销失效根源。",
    ]:
        add_paragraph(tf_left, item, font_size=17, color=PALETTE["ink"], bullet=True, space_after=8, line_spacing=1.15)

    add_round_rect(slide, 0.92, 6.04, 5.14, 0.44, PALETTE["mint"])
    add_textbox(
        slide,
        1.08,
        6.11,
        4.82,
        0.22,
        "研究目标：找出 FD 服饰短视频营销失效根源，并提出可落地的优化路径。",
        font_size=15,
        color=PALETTE["ink"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    add_round_rect(slide, 6.38, 2.74, 6.08, 1.78, PALETTE["white"], line=PALETTE["line"])
    add_panel_title(slide, 6.62, 2.94, "研究意义", size=18)
    add_textbox(slide, 6.66, 3.34, 1.18, 0.24, "理论层面", font_size=16, color=PALETTE["teal_dark"], bold=True)
    add_textbox(slide, 7.85, 3.3, 4.12, 0.5, "补充服饰品牌短视频营销案例研究，并验证 STP 与 4V 在细分赛道中的联动应用。", font_size=15, color=PALETTE["muted"])
    add_textbox(slide, 6.66, 3.9, 1.18, 0.24, "实践层面", font_size=16, color=PALETTE["teal_dark"], bold=True)
    add_textbox(slide, 7.85, 3.86, 4.12, 0.5, "为 FD 服饰提供从诊断到执行的短视频优化方案，也为同类中小服饰品牌提供可借鉴路径。", font_size=15, color=PALETTE["muted"])

    add_round_rect(slide, 6.38, 4.78, 6.08, 1.5, PALETTE["white"], line=PALETTE["line"])
    add_panel_title(slide, 6.62, 4.98, "研究目的与核心问题", size=18)
    add_textbox(slide, 6.72, 5.4, 5.5, 0.32, "FD 服饰的短视频问题，是战略层方向失准，还是执行层内容价值不足？", font_size=18, color=PALETTE["ink"], bold=True)
    add_textbox(slide, 6.72, 5.82, 5.38, 0.22, "论文以 STP 找方向，以 4V 找价值，最终落到优化策略与执行机制。", font_size=15, color=PALETTE["muted"])
    add_footer_bar(slide)


def build_literature_slide(prs: Presentation, slide_no: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_chrome(slide, "国内外研究现状", slide_no, total, subtitle="LITERATURE REVIEW")
    add_round_rect(slide, 0.6, 1.0, 12.1, 5.88, PALETTE["panel_soft"], line=PALETTE["line"])

    add_round_rect(slide, 0.86, 1.28, 5.55, 5.2, PALETTE["white"], line=PALETTE["line"])
    add_panel_title(slide, 1.06, 1.48, "国内研究", size=18)
    _, tf_cn = add_textbox(slide, 1.02, 1.92, 5.0, 4.2)
    for item in [
        "侯雪青（2024）指出，短视频的有用性与易用性会显著提升消费者购买意愿。",
        "刘雅婷、陈文晖（2025）提出“短视频种草 + 社交平台引流 + 直播带货转化”正在重塑服饰产业链。",
        "舒方瑞（2024）、严小林（2025）强调服饰短视频应突出版型、质感、细节与视觉修辞。",
        "现有国内研究更多聚焦头部品牌或行业层面，对中小服饰品牌的实操诊断仍偏少。",
    ]:
        add_paragraph(tf_cn, item, font_size=16, color=PALETTE["ink"], bullet=True, space_after=8, line_spacing=1.14)

    add_round_rect(slide, 6.68, 1.28, 5.56, 2.38, PALETTE["white"], line=PALETTE["line"])
    add_panel_title(slide, 6.9, 1.48, "国外研究", size=18)
    _, tf_foreign = add_textbox(slide, 6.86, 1.92, 4.96, 1.42)
    for item in [
        "Tam 等（2025）、Kanellos 等（2024）从数字营销与数据技术角度讨论品牌精细化触达。",
        "Luo 等（2025）与 Zafar 等（2021）认为内容沉浸、信任与互动机制会显著影响购买意愿。",
        "国外研究更重视前沿技术、全球趋势与消费者心理机制。",
    ]:
        add_paragraph(tf_foreign, item, font_size=15, color=PALETTE["ink"], bullet=True, space_after=8, line_spacing=1.14)

    add_round_rect(slide, 6.68, 3.98, 5.56, 2.5, PALETTE["white"], line=PALETTE["line"])
    add_panel_title(slide, 6.9, 4.18, "研究评述与论文切入点", size=18)
    _, tf_gap = add_textbox(slide, 6.86, 4.62, 4.96, 1.5)
    for item in [
        "现有研究较少聚焦 FD 这类中型服饰企业的真实运营问题。",
        "不少研究仅采用单一理论视角，难以同时解释“方向”与“价值”的双重问题。",
        "本文以 FD 为案例，构建 STP + 4V 双理论框架，提升诊断与策略建议的针对性。",
    ]:
        add_paragraph(tf_gap, item, font_size=15, color=PALETTE["ink"], bullet=True, space_after=8, line_spacing=1.14)
    add_footer_bar(slide)


def build_theory_slide(prs: Presentation, slide_no: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_chrome(slide, "理论依据", slide_no, total, subtitle="THEORETICAL BASIS")
    add_round_rect(slide, 0.6, 1.0, 12.1, 5.88, PALETTE["panel_soft"], line=PALETTE["line"])

    add_round_rect(slide, 0.86, 1.28, 3.64, 4.85, PALETTE["white"], line=PALETTE["line"])
    add_panel_title(slide, 1.08, 1.48, "STP 理论", size=18)
    add_textbox(slide, 1.12, 1.94, 2.8, 0.24, "Segmentation / Targeting / Positioning", font_size=12, color=PALETTE["muted"], font_name=FONT_EN)
    _, tf_stp = add_textbox(slide, 1.06, 2.34, 3.0, 2.95)
    for item in [
        "市场细分：按需求、场景、消费能力与价值偏好重新划分市场。",
        "目标市场：识别最值得服务的核心人群与次级人群。",
        "市场定位：明确品牌想被谁记住，以及凭什么被选择。",
        "作用：帮助论文回答“FD 应该做给谁看、如何被看见”。",
    ]:
        add_paragraph(tf_stp, item, font_size=15, color=PALETTE["ink"], bullet=True, space_after=7, line_spacing=1.13)

    add_round_rect(slide, 4.86, 1.28, 3.64, 4.85, PALETTE["white"], line=PALETTE["line"])
    add_panel_title(slide, 5.08, 1.48, "4V 理论", size=18)
    add_textbox(slide, 5.12, 1.94, 2.24, 0.24, "Variation / Versatility / Value / Vibration", font_size=12, color=PALETTE["muted"], font_name=FONT_EN)
    _, tf_4v = add_textbox(slide, 5.06, 2.34, 3.0, 2.95)
    for item in [
        "差异化：用鲜明内容形式与视觉标签提升识别度。",
        "功能化：让内容具备教程、决策辅助和转化引导作用。",
        "附加价值：输出品牌理念、设计故事与情绪价值。",
        "共鸣：通过互动、评论、UGC 和社群形成关系连接。",
    ]:
        add_paragraph(tf_4v, item, font_size=15, color=PALETTE["ink"], bullet=True, space_after=7, line_spacing=1.13)

    add_round_rect(slide, 8.86, 1.28, 3.38, 4.85, PALETTE["white"], line=PALETTE["line"])
    add_panel_title(slide, 9.08, 1.48, "双理论分析框架", size=18)
    chips = [
        ("1", "文献与行业背景"),
        ("2", "FD 现状与用户调研"),
        ("3", "STP 识别方向问题"),
        ("4", "4V 提出优化路径"),
    ]
    y = 2.02
    for idx, label in chips:
        add_chip(slide, 9.12, y, 2.72, 0.44, f"{idx}. {label}", fill=PALETTE["panel_soft"], color=PALETTE["teal_dark"], font_size=13)
        y += 0.62
    add_round_rect(slide, 9.1, 4.7, 2.84, 0.86, PALETTE["mint"])
    add_textbox(
        slide,
        9.22,
        4.84,
        2.58,
        0.5,
        "理论作用分工：STP 解决“方向”，4V 解决“价值”。",
        font_size=15,
        color=PALETTE["ink"],
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_VERTICAL_ANCHOR.MIDDLE,
    )
    add_footer_bar(slide)


def build_design_slide(prs: Presentation, slide_no: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_chrome(slide, "研究设计", slide_no, total, subtitle="RESEARCH DESIGN")
    add_round_rect(slide, 0.6, 1.0, 12.1, 5.75, PALETTE["panel_soft"], line=PALETTE["line"])

    add_round_rect(slide, 0.88, 1.38, 4.2, 5.0, PALETTE["white"], line=PALETTE["line"])
    add_textbox(slide, 1.08, 1.57, 1.6, 0.28, "研究方法", font_size=14, color=PALETTE["teal_dark"], bold=True)
    methods = [
        ("01", "文献分析法", "梳理短视频营销、服饰行业营销及 STP、4V 理论研究，搭建分析框架。"),
        ("02", "问卷调查法", "围绕短视频习惯、内容偏好、认知评价与购买意愿发放问卷，回收 412 份有效样本。"),
        ("03", "案例分析法", "结合 FD 服饰账号现状与内容表现，验证问题成因与策略可行性。"),
    ]
    y = 2.0
    for num, title, body in methods:
        add_round_rect(slide, 1.06, y, 0.6, 0.6, PALETTE["teal"])
        add_textbox(slide, 1.18, y + 0.09, 0.35, 0.2, num, font_size=13, color=PALETTE["white"], bold=True, font_name=FONT_EN, align=PP_ALIGN.CENTER)
        add_textbox(slide, 1.82, y + 0.03, 1.65, 0.22, title, font_size=15, color=PALETTE["ink"], bold=True)
        add_textbox(slide, 1.82, y + 0.32, 2.75, 0.52, body, font_size=11, color=PALETTE["muted"])
        y += 1.45

    add_round_rect(slide, 5.35, 1.38, 3.18, 2.2, PALETTE["white"], line=PALETTE["line"])
    add_textbox(slide, 5.56, 1.57, 2.0, 0.28, "调研样本与渠道", font_size=14, color=PALETTE["teal_dark"], bold=True)
    add_textbox(slide, 5.62, 1.95, 1.0, 0.45, "412", font_size=28, color=PALETTE["teal_dark"], bold=True, font_name=FONT_EN)
    add_textbox(slide, 6.58, 2.06, 1.3, 0.2, "份有效问卷", font_size=13, color=PALETTE["ink"], bold=True)
    _, tf_s = add_textbox(slide, 5.6, 2.48, 2.6, 0.9)
    add_paragraph(tf_s, "线上：时尚穿搭群、品牌粉丝群", font_size=13, color=PALETTE["muted"], bullet=True)
    add_paragraph(tf_s, "线下：FD 门店周边", font_size=13, color=PALETTE["muted"], bullet=True)

    add_round_rect(slide, 8.82, 1.38, 3.56, 2.2, PALETTE["white"], line=PALETTE["line"])
    add_textbox(slide, 9.05, 1.57, 2.0, 0.28, "研究逻辑", font_size=14, color=PALETTE["teal_dark"], bold=True)
    logic = ["行业背景与企业现状", "问卷调查与问题识别", "STP 诊断方向失准", "4V 优化执行路径"]
    x = 9.03
    for idx, label in enumerate(logic, start=1):
        add_round_rect(slide, x, 2.0, 0.52, 0.52, PALETTE["teal"], transparency=0.08)
        add_textbox(slide, x + 0.13, 2.12, 0.2, 0.15, str(idx), font_size=12, color=PALETTE["teal_dark"], bold=True, font_name=FONT_EN, align=PP_ALIGN.CENTER)
        add_textbox(slide, x - 0.03, 2.58, 0.95, 0.48, label, font_size=11, color=PALETTE["muted"], align=PP_ALIGN.CENTER)
        x += 0.83

    add_round_rect(slide, 5.35, 3.95, 7.03, 2.43, PALETTE["white"], line=PALETTE["line"])
    add_textbox(slide, 5.56, 4.16, 2.1, 0.28, "核心问题", font_size=14, color=PALETTE["teal_dark"], bold=True)
    add_textbox(
        slide,
        5.62,
        4.58,
        6.35,
        0.82,
        "FD 服饰的短视频营销问题，究竟出在战略层的方向失准，还是执行层的内容与运营偏差？",
        font_size=17,
        color=PALETTE["ink"],
        bold=True,
        valign=MSO_VERTICAL_ANCHOR.MIDDLE,
    )
    add_textbox(
        slide,
        5.62,
        5.45,
        6.25,
        0.42,
        "研究思路：用 STP 找准“该做给谁看、该如何定位”，再用 4V 重构“内容如何更有价值”。",
        font_size=13,
        color=PALETTE["muted"],
    )
    add_footer_bar(slide)


def build_status_slide(prs: Presentation, slide_no: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_chrome(slide, "企业现状与用户洞察", slide_no, total, subtitle="CURRENT STATUS")
    add_round_rect(slide, 0.6, 1.0, 12.1, 5.75, PALETTE["panel_soft"], line=PALETTE["line"])

    add_round_rect(slide, 0.88, 1.34, 4.2, 4.95, PALETTE["white"], line=PALETTE["line"])
    add_textbox(slide, 1.08, 1.55, 2.0, 0.28, "FD 服饰现状", font_size=14, color=PALETTE["teal_dark"], bold=True)
    _, tf = add_textbox(slide, 1.08, 1.92, 3.65, 1.42)
    for item in [
        "2013 年成立，本土时尚女装品牌。",
        "价格带集中在 150-500 元，面向都市年轻女性。",
        "已布局抖音、小红书等平台，但内容打法以产品展示为主。",
    ]:
        add_paragraph(tf, item, font_size=14, color=PALETTE["ink"], bullet=True, space_after=6)

    add_textbox(slide, 1.08, 3.65, 2.0, 0.26, "现有内容结构", font_size=14, color=PALETTE["teal_dark"], bold=True)
    bars = [
        (4.0, "产品上新展示", "60%", PALETTE["orange"]),
        (4.55, "日常 vlog / 直播切片", "30%", PALETTE["teal"]),
        (5.1, "穿搭教程", "10%", PALETTE["gold"]),
    ]
    for y, label, pct, color in bars:
        add_textbox(slide, 1.12, y, 1.6, 0.22, label, font_size=12, color=PALETTE["muted"])
        add_round_rect(slide, 2.4, y + 0.02, 1.9, 0.22, PALETTE["shadow"])
        width = {"60%": 1.52, "30%": 0.76, "10%": 0.25}[pct]
        add_round_rect(slide, 2.4, y + 0.02, width, 0.22, color)
        add_textbox(slide, 4.38, y - 0.02, 0.45, 0.25, pct, font_size=12, color=PALETTE["ink"], bold=True, align=PP_ALIGN.RIGHT)

    add_round_rect(slide, 5.35, 1.34, 7.0, 4.95, PALETTE["white"], line=PALETTE["line"])
    add_textbox(slide, 5.58, 1.55, 2.0, 0.28, "问卷关键发现", font_size=14, color=PALETTE["teal_dark"], bold=True)

    cards = [
        (5.58, 1.95, 1.42, "65%", "23-30 岁为主", 18),
        (7.12, 1.95, 1.42, "58%", "一二线城市", 18),
        (8.66, 1.95, 1.6, "3000-8000元", "月可支配收入", 14),
        (10.38, 1.95, 1.42, "1h+", "日均刷短视频时长", 18),
    ]
    for x, y, width, value, label, size in cards:
        add_round_rect(slide, x, y, width, 1.05, PALETTE["panel_soft"], line=PALETTE["line"])
        add_textbox(slide, x + 0.08, y + 0.16, width - 0.16, 0.28, value, font_size=size, color=PALETTE["teal_dark"], bold=True, font_name=FONT_EN, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.08, y + 0.56, width - 0.16, 0.28, label, font_size=10, color=PALETTE["muted"], align=PP_ALIGN.CENTER)

    add_round_rect(slide, 5.68, 3.38, 6.32, 1.0, PALETTE["mint"])
    add_textbox(slide, 5.86, 3.53, 5.95, 0.28, "内容偏好：72% 喜欢场景化穿搭教程，68% 关注真实测评，纯产品展示仅 35%。", font_size=13, color=PALETTE["ink"], bold=True)
    add_round_rect(slide, 5.68, 4.62, 6.32, 1.05, PALETTE["panel_soft"])
    add_textbox(slide, 5.86, 4.78, 5.95, 0.4, "认知评价：仅 28% 关注过 FD 短视频；超过八成认为内容同质化，75% 认为缺少搭配指导。", font_size=13, color=PALETTE["ink"], bold=True)
    add_footer_bar(slide)


def build_problem_slide(prs: Presentation, slide_no: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_chrome(slide, "问题诊断", slide_no, total, subtitle="DIAGNOSIS")
    add_round_rect(slide, 0.6, 1.0, 12.1, 5.75, PALETTE["panel_soft"], line=PALETTE["line"])

    add_round_rect(slide, 0.9, 1.42, 5.45, 4.75, PALETTE["white"], line=PALETTE["line"])
    add_textbox(slide, 1.12, 1.62, 2.6, 0.28, "STP 层面：方向失准", font_size=16, color=PALETTE["ink"], bold=True)
    stp_items = [
        ("市场细分粗放", "仍以“18-45 岁都市年轻女性”笼统覆盖，缺少需求与场景分层。"),
        ("目标市场摇摆", "核心客群与次级客群不清晰，内容锚点频繁变化。"),
        ("定位落地脱节", "“可持续时尚”没有进入内容体系，也没有形成固定栏目 IP。"),
    ]
    y = 2.12
    for idx, (title, body) in enumerate(stp_items, start=1):
        add_round_rect(slide, 1.12, y, 0.48, 0.48, PALETTE["teal"], transparency=0.08)
        add_textbox(slide, 1.22, y + 0.09, 0.18, 0.15, str(idx), font_size=12, color=PALETTE["teal_dark"], bold=True, font_name=FONT_EN, align=PP_ALIGN.CENTER)
        add_textbox(slide, 1.74, y - 0.02, 1.5, 0.22, title, font_size=14, color=PALETTE["teal_dark"], bold=True)
        add_textbox(slide, 1.72, y + 0.22, 4.15, 0.52, body, font_size=12, color=PALETTE["muted"])
        y += 1.15

    add_round_rect(slide, 6.62, 1.42, 5.45, 4.75, PALETTE["white"], line=PALETTE["line"])
    add_textbox(slide, 6.84, 1.62, 2.8, 0.28, "4V 层面：价值断裂", font_size=16, color=PALETTE["ink"], bold=True)
    four_v = [
        ("差异化不足", "内容形式单一，黄金 3 秒吸引力弱，账号风格不稳定。"),
        ("功能化不足", "教程、面料知识、尺码建议与 CTA 不完整。"),
        ("附加价值不足", "品牌故事、环保理念与会员权益没有被看见。"),
        ("共鸣不足", "缺少微剧情、UGC 互动与有温度的社群运营。"),
    ]
    y = 2.0
    for title, body in four_v:
        add_round_rect(slide, 6.9, y, 4.9, 0.82, PALETTE["panel_soft"])
        add_textbox(slide, 7.1, y + 0.1, 1.35, 0.2, title, font_size=13, color=PALETTE["teal_dark"], bold=True)
        add_textbox(slide, 8.3, y + 0.08, 3.2, 0.3, body, font_size=12, color=PALETTE["muted"])
        y += 0.95

    add_round_rect(slide, 0.9, 6.02, 11.17, 0.38, PALETTE["ink"])
    add_textbox(
        slide,
        1.1,
        6.07,
        10.75,
        0.2,
        "根源判断：FD 的短视频问题不是“发得不够多”，而是战略方向模糊、内容价值不强。",
        font_size=13,
        color=PALETTE["white"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer_bar(slide)


def build_stp_solution_slide(prs: Presentation, slide_no: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_chrome(slide, "基于 STP 的优化策略", slide_no, total, subtitle="STP SOLUTIONS")
    add_round_rect(slide, 0.6, 1.0, 12.1, 5.75, PALETTE["panel_soft"], line=PALETTE["line"])

    add_round_rect(slide, 0.88, 1.4, 5.55, 4.95, PALETTE["white"], line=PALETTE["line"])
    add_textbox(slide, 1.1, 1.6, 2.6, 0.28, "目标客群重构", font_size=14, color=PALETTE["teal_dark"], bold=True)
    segments = [
        ("核心目标", "精致通勤族", "22-28 岁，关注一衣多穿与通勤性价比。", PALETTE["orange"]),
        ("次级目标", "悦己体验党", "29-45 岁，更重视品质、设计与情绪价值。", PALETTE["teal"]),
        ("机会圈层", "可持续生活家", "关注环保面料、安全性与品牌理念一致性。", PALETTE["gold"]),
    ]
    y = 2.0
    for tag, title, body, color in segments:
        add_round_rect(slide, 1.1, y, 4.95, 0.92, PALETTE["panel_soft"])
        add_rect(slide, 1.1, y, 0.14, 0.92, color)
        add_textbox(slide, 1.35, y + 0.1, 1.0, 0.18, tag, font_size=10, color=color, bold=True)
        add_textbox(slide, 1.34, y + 0.28, 1.7, 0.22, title, font_size=15, color=PALETTE["ink"], bold=True)
        add_textbox(slide, 3.0, y + 0.18, 2.75, 0.45, body, font_size=12, color=PALETTE["muted"])
        y += 1.2

    add_round_rect(slide, 6.7, 1.4, 5.36, 4.95, PALETTE["white"], line=PALETTE["line"])
    add_textbox(slide, 6.95, 1.6, 2.6, 0.28, "内容结构优化", font_size=14, color=PALETTE["teal_dark"], bold=True)
    ratio_rows = [
        ("场景化通勤穿搭教程", "40%", PALETTE["orange"], 1.95),
        ("职场穿搭解决方案", "30%", PALETTE["teal"], 1.45),
        ("产品展示", "20%", PALETTE["gold"], 0.95),
        ("品牌故事 / 可持续理念", "10%", PALETTE["slate"], 0.45),
    ]
    y = 2.08
    for label, pct, color, width in ratio_rows:
        add_textbox(slide, 6.96, y - 0.02, 2.2, 0.22, label, font_size=12, color=PALETTE["muted"])
        add_round_rect(slide, 8.98, y, 2.3, 0.24, PALETTE["shadow"])
        add_round_rect(slide, 8.98, y, width, 0.24, color)
        add_textbox(slide, 11.38, y - 0.05, 0.42, 0.25, pct, font_size=11, color=PALETTE["ink"], bold=True, align=PP_ALIGN.RIGHT)
        y += 0.72

    add_round_rect(slide, 6.95, 5.1, 4.7, 0.78, PALETTE["mint"])
    add_textbox(
        slide,
        7.12,
        5.23,
        4.35,
        0.34,
        "账号锚点：围绕“FD 通勤穿搭指南”等系列栏目，持续强化品牌识别。",
        font_size=13,
        color=PALETTE["ink"],
        bold=True,
    )
    add_footer_bar(slide)


def build_4v_solution_slide(prs: Presentation, slide_no: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_chrome(slide, "基于 4V 的内容与运营策略", slide_no, total, subtitle="4V SOLUTIONS")
    add_round_rect(slide, 0.6, 1.0, 12.1, 5.75, PALETTE["panel_soft"], line=PALETTE["line"])

    cells = [
        (0.92, 1.4, "差异化", "真人出镜 + 痛点前置 + 黄金 3 秒，做出场景化栏目。"),
        (6.52, 1.4, "功能化", "一衣多穿、面料知识、尺码建议与明确 CTA 形成教程闭环。"),
        (0.92, 3.72, "附加价值", "品牌故事、设计幕后、环保实践与会员权益持续输出。"),
        (6.52, 3.72, "共鸣", "微剧情、UGC 征集、评论区答疑与社群运营一起发力。"),
    ]
    for x, y, title, body in cells:
        add_round_rect(slide, x, y, 5.04, 1.95, PALETTE["white"], line=PALETTE["line"])
        add_rect(slide, x, y, 0.12, 1.95, PALETTE["teal"])
        add_textbox(slide, x + 0.26, y + 0.18, 1.3, 0.24, title, font_size=17, color=PALETTE["ink"], bold=True)
        add_textbox(slide, x + 0.26, y + 0.56, 4.45, 0.52, body, font_size=13, color=PALETTE["muted"])

    add_textbox(slide, 1.18, 2.86, 4.3, 0.26, "典型动作：固定栏目、统一视觉规范、热点与品牌理念深度绑定。", font_size=11, color=PALETTE["muted"])
    add_textbox(slide, 6.78, 2.86, 4.2, 0.26, "典型动作：评论区问答、直播钩子、主页与购买链路一体优化。", font_size=11, color=PALETTE["muted"])
    add_textbox(slide, 1.18, 5.18, 4.2, 0.26, "典型动作：可持续叙事、设计师故事、知识内容与会员社群。", font_size=11, color=PALETTE["muted"])
    add_textbox(slide, 6.78, 5.18, 4.25, 0.26, "典型动作：微剧情、穿搭挑战、UGC 征集与精细化互动回复。", font_size=11, color=PALETTE["muted"])

    add_round_rect(slide, 2.2, 6.05, 8.9, 0.36, PALETTE["ink"])
    add_textbox(
        slide,
        2.38,
        6.09,
        8.55,
        0.2,
        "目标结果：让 FD 从“展示商品”转向“用内容创造价值、用品牌沉淀复购”。",
        font_size=13,
        color=PALETTE["white"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer_bar(slide)


def build_conclusion_slide(prs: Presentation, slide_no: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_chrome(slide, "研究结论与展望", slide_no, total, subtitle="CONCLUSION")
    add_round_rect(slide, 0.6, 1.0, 12.1, 5.75, PALETTE["panel_soft"], line=PALETTE["line"])

    blocks = [
        (
            0.92,
            "研究结论",
            [
                "FD 的问题不在“视频拍得太少”，而在战略失焦与执行失衡并存。",
                "优化逻辑应先以 STP 校准方向，再用 4V 提升内容价值与用户连接。",
                "这套路径对中小服饰品牌同样具有参考意义。",
            ],
        ),
        (
            4.36,
            "研究局限",
            [
                "研究对象为单案例，结论更贴合中小服饰企业。",
                "样本以一二线年轻女性为主，对下沉市场覆盖不足。",
                "策略仍停留在理论层面，尚缺长期市场验证。",
            ],
        ),
        (
            7.8,
            "未来展望",
            [
                "借助 AI 做内容洞察、文案生成与虚拟试穿提效。",
                "搭建“头部 + 腰部 + 素人”的金字塔达人合作矩阵。",
                "优化组织架构与考核体系，形成内容驱动的增长闭环。",
            ],
        ),
    ]
    for x, title, items in blocks:
        add_round_rect(slide, x, 1.45, 3.05, 4.7, PALETTE["white"], line=PALETTE["line"])
        add_textbox(slide, x + 0.2, 1.67, 1.5, 0.3, title, font_size=16, color=PALETTE["ink"], bold=True)
        _, tf = add_textbox(slide, x + 0.18, 2.1, 2.55, 3.4)
        for item in items:
            add_paragraph(tf, item, font_size=13, color=PALETTE["muted"], bullet=True, space_after=7)

    add_round_rect(slide, 0.92, 5.63, 9.93, 0.42, PALETTE["mint"])
    add_textbox(
        slide,
        1.1,
        5.69,
        9.55,
        0.2,
        "一句话总结：FD 服饰要真正破局，关键不是追热点，而是持续提供被目标用户需要的内容价值。",
        font_size=12,
        color=PALETTE["ink"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer_bar(slide)


def build_thanks_slide(prs: Presentation, slide_no: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, PALETTE["teal_deep"])
    add_rect(slide, 0, 0, 13.333, 7.5, PALETTE["teal_deep"], transparency=0.04)
    add_rect(slide, 8.22, 0.98, 4.5, 5.35, PALETTE["white"], transparency=0.1)
    add_picture_fill(slide, IMAGES["thanks"], 8.02, 1.03, 4.4, 5.3)
    if IMAGES["logo"].exists():
        add_picture_fill(slide, IMAGES["logo"], 0.5, 0.22, 1.35, 0.32)
    add_textbox(slide, 0.72, 1.62, 3.2, 0.6, "THANKS!", font_size=28, color=PALETTE["white"], bold=True, font_name=FONT_EN)
    add_textbox(slide, 0.74, 3.02, 3.0, 0.52, "谢谢聆听", font_size=24, color=PALETTE["white"], bold=True)
    add_textbox(slide, 0.76, 3.72, 4.0, 0.42, "恳请各位老师批评指正", font_size=16, color=PALETTE["white"])
    add_textbox(slide, 0.76, 5.55, 4.75, 0.38, "FD服饰公司短视频营销策略优化研究", font_size=14, color=PALETTE["white"], bold=True)
    add_textbox(slide, 0.76, 6.1, 4.7, 0.32, "答辩人：廖子楚    指导教师：巨静文", font_size=13, color=PALETTE["white"])
    add_textbox(slide, 12.12, 0.18, 0.35, 0.22, f"{slide_no:02d}", font_size=10, color=PALETTE["white"], font_name=FONT_EN)


def build_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides = [
        lambda: build_cover_slide(prs, 1, 15),
        lambda: build_contents_slide(prs, 2, 15),
        lambda: build_section_slide(prs, 3, 15, "01", "研究背景与意义", IMAGES["section1"]),
        lambda: build_background_slide(prs, 4, 15),
        lambda: build_section_slide(prs, 5, 15, "02", "研究设计", IMAGES["section2"]),
        lambda: build_design_slide(prs, 6, 15),
        lambda: build_section_slide(prs, 7, 15, "03", "现状与问题诊断", IMAGES["section3"]),
        lambda: build_status_slide(prs, 8, 15),
        lambda: build_problem_slide(prs, 9, 15),
        lambda: build_section_slide(prs, 10, 15, "04", "优化策略", IMAGES["section4"]),
        lambda: build_stp_solution_slide(prs, 11, 15),
        lambda: build_4v_solution_slide(prs, 12, 15),
        lambda: build_section_slide(prs, 13, 15, "05", "研究结论与展望", IMAGES["section5"]),
        lambda: build_conclusion_slide(prs, 14, 15),
        lambda: build_thanks_slide(prs, 15, 15),
    ]

    for make in slides:
        make()

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    return OUTPUT_PATH


if __name__ == "__main__":
    path = build_presentation()
    print(path)
